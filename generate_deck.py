"""
generate_deck.py
Generates the Trase × RxCo Product Demo & Impact Analysis presentation.

Usage:
    pip install -r requirements.txt
    python generate_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MID_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x00, 0xAA, 0x44)

OUTPUT_FILE = "Trase_x_RxCo_Product_Demo_and_Impact_Analysis.pptx"

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6]  # completely blank


def fill_solid(shape, colour: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, colour: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    fill_solid(shape, colour)
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_font(run, size_pt, bold=False, colour: RGBColor = None, font_name="Calibri"):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if colour:
        run.font.color.rgb = colour


def add_text_to_box(txBox, text, size_pt, bold=False, colour: RGBColor = WHITE,
                    alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, colour, font_name)
    return tf


def add_paragraph(tf, text, size_pt, bold=False, colour: RGBColor = DARK_GRAY,
                  alignment=PP_ALIGN.LEFT, space_before_pt=0, font_name="Calibri"):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = _Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, colour, font_name)
    return p


def add_footer(slide, text="TRASE | Confidential"):
    tb = add_textbox(slide, Inches(0.3), Inches(7.1), Inches(4), Inches(0.3))
    add_text_to_box(tb, text, 9, False, MID_GRAY)


def bg_rect(slide, colour: RGBColor = NAVY):
    """Full-bleed background rectangle."""
    return add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colour)


def add_title_bar(slide, title_text, size_pt=28):
    """Dark navy title bar across the top of a content slide."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), NAVY)
    tb = add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.7))
    add_text_to_box(tb, title_text, size_pt, True, WHITE)
    return bar


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def add_table(slide, rows, cols, left, top, width, height,
              headers=None, data=None,
              header_bg=NAVY, header_fg=WHITE,
              alt_row_bg=LIGHT_GRAY,
              col_widths=None,
              font_size=12):
    """Add a styled table to a slide."""
    from pptx.util import Pt as _Pt

    table_rows = (1 if headers else 0) + (len(data) if data else 0)
    tbl = slide.shapes.add_table(table_rows, cols, left, top, width, height).table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)

    def _set_cell(cell, text, bg, fg, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = fg

    row_idx = 0
    if headers:
        for c, h in enumerate(headers):
            _set_cell(tbl.cell(0, c), h, header_bg, header_fg, bold=True)
        row_idx = 1

    if data:
        for r, row_data in enumerate(data):
            bg = alt_row_bg if r % 2 == 0 else WHITE
            for c, cell_text in enumerate(row_data):
                _set_cell(tbl.cell(row_idx + r, c), cell_text, bg, DARK_GRAY)

    return tbl


# ---------------------------------------------------------------------------
# KPI callout box
# ---------------------------------------------------------------------------

def add_kpi_box(slide, left, top, width, height, lines, bg=TEAL):
    rect = add_rect(slide, left, top, width, height, bg)
    rect.line.fill.background()
    tb = add_textbox(slide, left + Inches(0.1), top + Inches(0.1),
                     width - Inches(0.2), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_font(run, size, bold, WHITE)


# ===========================================================================
# SLIDES
# ===========================================================================

def slide_01_title(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    bg_rect(slide, NAVY)

    # Accent bar at bottom
    add_rect(slide, 0, Inches(6.8), SLIDE_W, Inches(0.7), TEAL)

    # "TRASE" logo text top-left
    tb = add_textbox(slide, Inches(0.4), Inches(0.25), Inches(3), Inches(0.6))
    add_text_to_box(tb, "TRASE", 22, True, TEAL)

    # Title
    tb = add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
    add_text_to_box(tb, "Trase \u00d7 RxCo", 54, True, WHITE, PP_ALIGN.CENTER)

    # Subtitle
    tb = add_textbox(slide, Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.9))
    add_text_to_box(tb, "Product Demo & Impact Analysis", 30, False, TEAL, PP_ALIGN.CENTER)

    # Footer
    tb = add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5))
    add_text_to_box(
        tb,
        "Prepared for: RxCo Leadership Team  |  Date: March 16, 2026  |  "
        "Prepared by: Technical Program Manager, Trase",
        11, False, WHITE, PP_ALIGN.CENTER,
    )


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Today's Agenda")

    items = [
        "1.  RxCo Workflow Overview",
        "2.  Product Demo \u2014 Eligibility Determination Agent\n"
        "     \u2022 Live walkthrough    \u2022 Technical architecture    \u2022 Edge case handling",
        "3.  Impact Analysis\n"
        "     \u2022 Current cost baseline    \u2022 Savings by automation tier    \u2022 12-month ROI projection",
        "4.  Assumptions & Validation Plan",
        "5.  Recommended Next Steps",
    ]

    tb = add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = item
        set_font(run, 16, False, DARK_GRAY)

    add_footer(slide)


def slide_03_exec_summary(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Executive Summary")

    body = (
        "RxCo\u2019s enrollment workflow currently costs $1.50M\u2013$2.00M annually at the "
        "50,000-member scale. Trase\u2019s AI agents can automate the most labor-intensive, "
        "rules-based steps \u2014 starting with Eligibility Determination \u2014 and progressively "
        "extend across the full six-step workflow."
    )
    tb = add_textbox(slide, Inches(0.4), Inches(1.0), Inches(8.5), Inches(2.5))
    add_text_to_box(tb, body, 16, False, DARK_GRAY)

    add_kpi_box(
        slide,
        Inches(9.2), Inches(1.0), Inches(3.8), Inches(5.8),
        [
            ("3.1\u00d7", 44, True),
            ("Projected Year 1 ROI", 14, False),
            ("", 10, False),
            ("$599K", 36, True),
            ("Net Annual Savings", 14, False),
            ("", 10, False),
            ("Exceeds 2\u00d7 guarantee by 55%", 13, True),
        ],
    )

    add_footer(slide)


def slide_04_business_model(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Business Model Overview")

    add_table(
        slide,
        rows=3, cols=3,
        left=Inches(0.4), top=Inches(1.1),
        width=Inches(12.5), height=Inches(1.6),
        headers=["Channel", "Description", "Pricing"],
        data=[
            ["D2C", "Individual patients enroll directly", "$50 / member / month / Rx"],
            ["B2B", "MSOs, self-funded orgs, employer groups", "15% shared-savings model"],
        ],
        col_widths=[2, 6, 4],
        font_size=14,
    )

    tb = add_textbox(slide, Inches(0.4), Inches(3.0), Inches(12.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    metrics = [
        ("Key Metrics", 18, True, NAVY),
        ("\u2022  50,000 enrolled members  (25K D2C / 25K B2B)", 15, False, DARK_GRAY),
        ("\u2022  Average medication cost per patient: $2,500", 15, False, DARK_GRAY),
        ("\u2022  RxCo revenue per B2B case: $375", 15, False, DARK_GRAY),
    ]
    first = True
    for text, size, bold, colour in metrics:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = text
        set_font(run, size, bold, colour)

    add_footer(slide)


def slide_05_workflow(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(
        slide,
        "Current Enrollment Workflow \u2014 6 Steps, 90\u2013120 Minutes, $30\u2013$40 per Case",
        size_pt=22,
    )

    steps = [
        ("[1]\nEligibility\nDetermination", "28 min | $9.33"),
        ("[2]\nPatient\nContact", "18 min | $6.00"),
        ("[3]\nPatient\nAuthorization", "12 min | $4.00"),
        ("[4]\nProvider\nCoordination", "17 min | $5.67"),
        ("[5]\nPharma\nSubmission", "18 min | $6.00"),
        ("[6]\nFollow-Up &\nRenewal", "12 min | $4.00"),
    ]

    box_w = Inches(1.95)
    box_h = Inches(2.0)
    top = Inches(1.1)
    gap = Inches(0.12)
    start_left = Inches(0.25)

    for i, (label, detail) in enumerate(steps):
        left = start_left + i * (box_w + gap)
        rect = add_rect(slide, left, top, box_w, box_h, NAVY)

        tb = add_textbox(slide, left + Inches(0.1), top + Inches(0.1),
                         box_w - Inches(0.2), box_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_font(run, 12, True, WHITE)

        tb2 = add_textbox(slide, left + Inches(0.05), top + box_h - Inches(0.45),
                          box_w - Inches(0.1), Inches(0.4))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = detail
        set_font(run2, 11, False, TEAL)

    # Callout at bottom
    add_kpi_box(
        slide,
        Inches(3.0), Inches(5.5), Inches(7.3), Inches(0.8),
        [("Annual cost at 50,000 members: $1,750,000", 16, True)],
    )

    add_footer(slide)


def slide_06_automation_roadmap(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Phased Automation Roadmap")

    tiers = [
        (
            "TIER 1 \u2014 Immediate (Pre-Launch)",
            "Eligibility Determination  |  No patient interaction  |  Deploy now",
            TEAL,
        ),
        (
            "TIER 2 \u2014 Near-Term (Post-Launch)",
            "Document processing, form population, signature collection, pharma submission",
            RGBColor(0x00, 0x99, 0xBB),
        ),
        (
            "TIER 3 \u2014 Deferred",
            "Patient interaction agents  |  Provider review optimization  |  Benchmark against call center first",
            RGBColor(0x00, 0x66, 0x88),
        ),
    ]

    top = Inches(1.1)
    for i, (heading, detail, colour) in enumerate(tiers):
        band_top = top + i * Inches(1.8)
        add_rect(slide, Inches(0.3), band_top, Inches(12.7), Inches(1.6), colour)

        tb = add_textbox(slide, Inches(0.5), band_top + Inches(0.1),
                         Inches(12.3), Inches(0.55))
        add_text_to_box(tb, heading, 18, True, WHITE)

        tb2 = add_textbox(slide, Inches(0.5), band_top + Inches(0.65),
                          Inches(12.3), Inches(0.8))
        add_text_to_box(tb2, detail, 15, False, WHITE)

    add_footer(slide)


def slide_07_divider_demo(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    bg_rect(slide, NAVY)

    add_rect(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.06), TEAL)

    tb = add_textbox(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2))
    add_text_to_box(tb, "PRODUCT DEMO", 48, True, WHITE, PP_ALIGN.CENTER)

    tb2 = add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9))
    add_text_to_box(tb2, "Eligibility Determination Agent", 26, False, TEAL, PP_ALIGN.CENTER)

    tb3 = add_textbox(slide, Inches(0.3), Inches(0.2), Inches(2), Inches(0.5))
    add_text_to_box(tb3, "TRASE", 18, True, TEAL)


def slide_08_why_elig(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Why We\u2019re Demoing Step 1: Eligibility Determination")

    reasons = [
        "1.  Agreed starting point \u2014 aligned in pilot scoping",
        "2.  Zero patient interaction risk \u2014 operates on claims data only",
        "3.  Highest time consumption \u2014 28 min (25% of total workflow)",
        "4.  Rules-based & auditable \u2014 ideal for healthcare compliance",
        "5.  Volume multiplier \u2014 scales worst manually, best automated",
    ]

    tb = add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for r in reasons:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(14)
        run = p.add_run()
        run.text = r
        set_font(run, 17, False, DARK_GRAY)

    add_footer(slide)


def slide_09_demo_scenario(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Demo Flow \u2014 What You\u2019ll See")

    phases = [
        ("Phase 1: UPLOAD", "Operator uploads 500-record synthetic claims file"),
        ("Phase 2: CONFIGURE", "Operator reviews / adjusts eligibility rules"),
        ("Phase 3: PROCESS", "Agent parses, cross-references, scores (45 sec)"),
        ("Phase 4: RESULTS", "Dashboard: ranked patients, scores, exceptions"),
        ("Phase 5: REVIEW", "Operator inspects 3 flagged edge cases"),
        ("Phase 6: EXPORT", 'CSV + PDF export, "Ready for Step 2" queue'),
    ]

    box_w = Inches(1.95)
    box_h = Inches(1.5)
    top = Inches(1.1)
    gap = Inches(0.12)
    start_left = Inches(0.25)

    for i, (phase, detail) in enumerate(phases):
        left = start_left + i * (box_w + gap)
        add_rect(slide, left, top, box_w, box_h, NAVY)

        tb = add_textbox(slide, left + Inches(0.05), top + Inches(0.1),
                         box_w - Inches(0.1), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = phase
        set_font(run, 11, True, TEAL)

        tb2 = add_textbox(slide, left + Inches(0.05), top + Inches(0.65),
                          box_w - Inches(0.1), Inches(0.8))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = detail
        set_font(run2, 11, False, WHITE)

    stats = (
        "Results: 412 eligible  |  27 flagged  |  61 ineligible  |  "
        "Processing: 45 seconds vs. ~210 hours manually"
    )
    add_kpi_box(
        slide,
        Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.8),
        [(stats, 14, True)],
    )

    add_footer(slide)


def slide_10_architecture(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Technical Architecture")

    flow = [
        ("Claims File\nUpload\n(CSV/Excel/HL7)", WHITE, NAVY),
        ("Ingestion &\nParsing Engine\n[LIVE]", WHITE, NAVY),
        ("Rules Engine\n[LIVE]\n\u2022 LIS check\n\u2022 Insurance filter\n\u2022 Copay threshold\n\u2022 Formulary match",
         WHITE, NAVY),
        ("Drug Formulary\nDB\n[MOCK \u2014 10 meds]", WHITE, RGBColor(0x99, 0x66, 0x00)),
        ("Eligibility\nScoring &\nRanking [LIVE]", WHITE, NAVY),
        ("Output\nDashboard\n[MOCK] &\nExport API\n[LIVE]", WHITE, NAVY),
    ]

    box_w = Inches(1.95)
    box_h = Inches(2.5)
    top = Inches(1.1)
    gap = Inches(0.1)
    start_left = Inches(0.2)

    for i, (label, fg, bg) in enumerate(flow):
        left = start_left + i * (box_w + gap)
        add_rect(slide, left, top, box_w, box_h, bg)
        tb = add_textbox(slide, left + Inches(0.05), top + Inches(0.15),
                         box_w - Inches(0.1), box_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_font(run, 12, False, fg)

        if i < len(flow) - 1:
            arrow_left = left + box_w + Inches(0.01)
            tb_arr = add_textbox(slide, arrow_left, top + Inches(1.1), Inches(0.08), Inches(0.3))
            add_text_to_box(tb_arr, "\u2192", 16, True, DARK_GRAY)

    # Legend
    add_kpi_box(slide, Inches(0.3), Inches(5.0), Inches(5.5), Inches(0.7),
                [("Green = Live Production  |  Yellow = Mock / Prototype", 12, False)],
                bg=RGBColor(0x22, 0x66, 0x44))

    add_footer(slide)


def slide_11_edge_cases(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "How the Agent Handles Edge Cases")

    add_table(
        slide,
        rows=7, cols=3,
        left=Inches(0.3), top=Inches(1.0),
        width=Inches(12.7), height=Inches(5.8),
        headers=["Scenario", "Agent Response", "Trust Message"],
        data=[
            ["Malformed data",
             "Flags errors, processes clean records",
             "\u201cNothing fails silently\u201d"],
            ["Multi-program match",
             "Ranks by savings + approval likelihood",
             "\u201cBest option with full reasoning\u201d"],
            ["Low confidence (<70)",
             "Routes to exception queue with analysis",
             "\u201cPre-analyzed ambiguity\u201d"],
            ["System latency",
             "Shows partial results, continues async",
             "\u201cNothing is lost\u201d"],
            ["Borderline threshold",
             "Flags as near-threshold, one-click override",
             "\u201cVolume for agent, judgment for you\u201d"],
            ["Drug not in formulary",
             "Logs as coverage gap insight",
             "\u201cStrategic intelligence\u201d"],
        ],
        col_widths=[3, 5, 4],
        font_size=13,
    )

    add_footer(slide)


def slide_12_divider_impact(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    bg_rect(slide, NAVY)

    add_rect(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.06), TEAL)

    tb = add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2))
    add_text_to_box(tb, "IMPACT ANALYSIS", 48, True, WHITE, PP_ALIGN.CENTER)

    tb2 = add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9))
    add_text_to_box(
        tb2,
        "Quantifying the Opportunity Across All 6 Steps",
        24, False, TEAL, PP_ALIGN.CENTER,
    )

    tb3 = add_textbox(slide, Inches(0.3), Inches(0.2), Inches(2), Inches(0.5))
    add_text_to_box(tb3, "TRASE", 18, True, TEAL)


def slide_13_cost_baseline(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Current Manual Enrollment Cost")

    add_table(
        slide,
        rows=8, cols=3,
        left=Inches(0.4), top=Inches(1.0),
        width=Inches(9.0), height=Inches(5.5),
        headers=["Step", "Time (min)", "Cost / Case"],
        data=[
            ["1. Eligibility Determination", "28", "$9.33"],
            ["2. Initial Patient Contact", "18", "$6.00"],
            ["3. Patient Authorization", "12", "$4.00"],
            ["4. Provider Coordination", "17", "$5.67"],
            ["5. Pharma Submission", "18", "$6.00"],
            ["6. Follow-Up & Renewal", "12", "$4.00"],
            ["TOTAL", "105", "$35.00"],
        ],
        col_widths=[6, 2, 2],
        font_size=14,
    )

    add_kpi_box(
        slide,
        Inches(9.6), Inches(1.0), Inches(3.5), Inches(2.2),
        [
            ("$1,750,000", 36, True),
            ("Annual Cost at", 14, False),
            ("50,000 Members", 14, False),
        ],
    )

    add_footer(slide)


def slide_14_savings(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Annual Savings by Tier (50,000 Members)")

    tiers = [
        (
            "TIER 1 \u2014 $396,667",
            "Step 1: 90% automation rate, 85% time reduction",
            TEAL,
        ),
        (
            "TIER 2 \u2014 $540,000",
            "Step 3: $160K  |  Step 4: $170K  |  Step 5: $210K",
            RGBColor(0x00, 0x99, 0xBB),
        ),
        (
            "TIER 3 \u2014 $230,000",
            "Step 2: $120K  |  Step 6: $110K",
            RGBColor(0x00, 0x66, 0x88),
        ),
    ]

    top = Inches(1.1)
    for i, (heading, detail, colour) in enumerate(tiers):
        band_top = top + i * Inches(1.7)
        add_rect(slide, Inches(0.3), band_top, Inches(12.7), Inches(1.5), colour)
        tb = add_textbox(slide, Inches(0.5), band_top + Inches(0.1),
                         Inches(12.3), Inches(0.55))
        add_text_to_box(tb, heading, 20, True, WHITE)
        tb2 = add_textbox(slide, Inches(0.5), band_top + Inches(0.65),
                          Inches(12.3), Inches(0.7))
        add_text_to_box(tb2, detail, 15, False, WHITE)

    add_kpi_box(
        slide,
        Inches(3.5), Inches(6.25), Inches(6.3), Inches(0.85),
        [("TOTAL GROSS SAVINGS: $1,166,667", 20, True)],
    )

    add_footer(slide)


def slide_15_roi(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "12-Month ROI Projection \u2014 All Tiers Combined")

    add_table(
        slide,
        rows=7, cols=2,
        left=Inches(0.4), top=Inches(1.0),
        width=Inches(8.0), height=Inches(4.8),
        headers=["Metric", "Value"],
        data=[
            ["Total Gross Savings", "$892,122"],
            ["Trase Fees", "$275,125"],
            ["HITL Overhead", "$17,917"],
            ["NET SAVINGS", "$599,081"],
            ["ROI", "3.1\u00d7"],
            ["Guarantee met?", "\u2705 Yes"],
        ],
        col_widths=[5, 3],
        font_size=14,
    )

    tb = add_textbox(slide, Inches(0.4), Inches(6.0), Inches(8.0), Inches(0.4))
    add_text_to_box(
        tb,
        "Ramp-Up: 70% automation Month 1 \u2192 92% by Month 4+  |  "
        "Accuracy: 88% Month 1 \u2192 97% Month 4+",
        12, False, DARK_GRAY,
    )

    add_kpi_box(
        slide,
        Inches(8.8), Inches(1.0), Inches(4.2), Inches(4.8),
        [
            ("3.1\u00d7 ROI", 44, True),
            ("", 8, False),
            ("$599K Net Savings", 20, True),
            ("", 8, False),
            ("Guarantee: 2\u00d7", 16, False),
            ("Delivered: 3.1\u00d7", 16, False),
            ("55% above target", 16, True),
        ],
    )

    add_footer(slide)


def slide_16_what_it_means(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "What This Means for RxCo")

    kpis = [
        ("23,300\nNURSE-HOURS FREED",
         "Equivalent of 11.2 FTEs redirected from\ndata mining to patient engagement"),
        ("57\u201365%\nCOST REDUCTION",
         "Marginal enrollment cost drops from\n$35 \u2192 $12\u2013$15 / member"),
        ("$525K\nMARGIN RETAINED",
         "B2B enrollment cost drag drops from\n9.3% \u2192 3.7%"),
        ("28 MIN \u2192 5 SEC",
         "Eligibility determination processing time\n(+ 4 min human review)"),
    ]

    box_w = Inches(5.8)
    box_h = Inches(2.5)
    positions = [
        (Inches(0.4), Inches(1.0)),
        (Inches(6.9), Inches(1.0)),
        (Inches(0.4), Inches(3.85)),
        (Inches(6.9), Inches(3.85)),
    ]

    for (left, top), (headline, detail) in zip(positions, kpis):
        add_rect(slide, left, top, box_w, box_h, NAVY)

        tb = add_textbox(slide, left + Inches(0.15), top + Inches(0.15),
                         box_w - Inches(0.3), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = headline
        set_font(run, 20, True, TEAL)

        tb2 = add_textbox(slide, left + Inches(0.15), top + Inches(1.3),
                          box_w - Inches(0.3), Inches(1.1))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = detail
        set_font(run2, 13, False, WHITE)

    add_footer(slide)


def slide_17_assumptions(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Key Assumptions & How We\u2019ll Validate in the Pilot")

    add_table(
        slide,
        rows=7, cols=2,
        left=Inches(0.3), top=Inches(1.0),
        width=Inches(12.7), height=Inches(5.8),
        headers=["Assumption", "Pilot Validation"],
        data=[
            ["$20/hr labor rate",
             "Confirm actual blended rate with RxCo Ops"],
            ["105 min / enrollment",
             "Time-study: shadow 50 cases in Weeks 1\u20132"],
            ["90% automation (Tier 1)",
             "Track daily; target \u226585% by Week 6"],
            ["88%\u219297% accuracy ramp",
             "Audit 20 decisions/day; retrain weekly"],
            ["50K annual volume",
             "Confirm run-rate; model at 40K & 75K"],
            ["$3\u2192$2 HITL cost",
             "Measure nurse review time; should decrease"],
        ],
        col_widths=[5, 7],
        font_size=14,
    )

    add_footer(slide)


def slide_18_next_steps(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title_bar(slide, "Recommended Next Steps")

    timeline = [
        ("March 19", "RxCo provides claims schema + eligibility criteria"),
        ("March 19\u201326", "Trase builds demo (1 week)"),
        ("March 27", "Demo walkthrough with RxCo leadership"),
        ("March 28", "8-week pilot begins (1 practice, 3\u20134 locations)"),
        ("Fridays", "Weekly accuracy & performance reviews"),
        ("May 22", "Pilot impact card delivered"),
        ("May 25\u201329", "Agree on long-term subscription rate"),
        ("June 2026", "Scope Tier 2 automations for Q3 deployment"),
    ]

    tb = add_textbox(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for date, action in timeline:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run_date = p.add_run()
        run_date.text = f"{date}: "
        set_font(run_date, 15, True, NAVY)
        run_action = p.add_run()
        run_action.text = action
        set_font(run_action, 15, False, DARK_GRAY)

    # Guarantee callout
    add_kpi_box(
        slide,
        Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.8),
        [
            ("Trase\u2019s 2\u00d7 ROI guarantee applies contractually \u2014 "
             "if we don\u2019t deliver, you don\u2019t pay.",
             14, True)
        ],
    )

    add_footer(slide)


# ===========================================================================
# Main
# ===========================================================================

def build_deck():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_exec_summary(prs)
    slide_04_business_model(prs)
    slide_05_workflow(prs)
    slide_06_automation_roadmap(prs)
    slide_07_divider_demo(prs)
    slide_08_why_elig(prs)
    slide_09_demo_scenario(prs)
    slide_10_architecture(prs)
    slide_11_edge_cases(prs)
    slide_12_divider_impact(prs)
    slide_13_cost_baseline(prs)
    slide_14_savings(prs)
    slide_15_roi(prs)
    slide_16_what_it_means(prs)
    slide_17_assumptions(prs)
    slide_18_next_steps(prs)

    prs.save(OUTPUT_FILE)
    print(f"Saved: {OUTPUT_FILE}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_deck()
